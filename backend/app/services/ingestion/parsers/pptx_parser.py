from io import BytesIO

def parse_pptx(data: bytes) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX ingestion requires the python-pptx dependency") from exc
    presentation = Presentation(BytesIO(data))
    documents = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shapes = sorted(slide.shapes, key=lambda shape: (getattr(shape, "top", 0), getattr(shape, "left", 0)))
        texts = [shape.text.strip() for shape in shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        if not texts:
            continue
        title = getattr(slide.shapes.title, "text", "").strip() if slide.shapes.title else ""
        body = "\n".join(texts)
        if title and title not in body:
            body = f"{title}\n{body}"
        documents.append({"text": body, "metadata": {"source_locator": {"type": "slide", "value": slide_number}}})
    return documents
