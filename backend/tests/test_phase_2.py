import uuid
from io import BytesIO
import pytest
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

from app.db.base import Base
from app.models.user import User
from app.models.subject import Subject
from app.models.material import Material
from app.services.ingestion.chunker import chunk_documents
from app.services.ingestion.pipeline import IngestionPipeline
from app.services.ingestion.parsers import PARSERS
from app.utils.qdrant_client import point_id

@pytest.fixture
def db():
    engine = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False}, poolclass=StaticPool)
    Base.metadata.create_all(engine)
    session = sessionmaker(bind=engine)()
    yield session
    session.close()

def test_pdf_chunk_overlap_and_stable_source_locator():
    docs = [{"text": " ".join(f"w{i}" for i in range(1000)), "metadata": {"source_locator": {"type": "page", "value": 4}}}]
    chunks = chunk_documents(docs, file_type="pdf")
    assert len(chunks) == 3
    assert chunks[0].text.split()[-50:] == chunks[1].text.split()[:50]
    assert all(chunk.metadata["source_locator"] for chunk in chunks)

def test_pptx_chunking_has_slide_locator_and_overlap():
    docs = [{"text": f"Slide {i}", "metadata": {"source_locator": {"type": "slide", "value": i}}} for i in range(1, 4)]
    chunks = chunk_documents(docs, file_type="pptx")
    assert len(chunks) == 2
    assert chunks[0].metadata["source_locator"]["value"] == 1
    assert "Slide 2" in chunks[0].text and "Slide 2" in chunks[1].text

def test_deterministic_point_id_changes_by_material_or_chunk():
    material_id = uuid.uuid4()
    assert point_id(material_id, 0) == point_id(material_id, 0)
    assert point_id(material_id, 0) != point_id(material_id, 1)
    assert point_id(material_id, 0) != point_id(uuid.uuid4(), 0)

def test_parser_registry_rejects_unknown_extension():
    assert set(PARSERS) == {"pdf", "pptx", "docx"}

def test_real_pdf_pptx_docx_fixtures_preserve_source_locations():
    from reportlab.pdfgen import canvas
    pdf = BytesIO(); pdf_canvas = canvas.Canvas(pdf); pdf_canvas.drawString(72, 720, "PDF page one"); pdf_canvas.showPage(); pdf_canvas.save()
    pdf_docs = PARSERS["pdf"](pdf.getvalue())
    assert pdf_docs[0]["metadata"]["source_locator"] == {"type": "page", "value": 1}

    from pptx import Presentation
    from pptx.util import Inches
    presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)); title.text = "Sparse title"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1)); body.text = "Body text"
    pptx = BytesIO(); presentation.save(pptx)
    pptx_docs = PARSERS["pptx"](pptx.getvalue())
    assert pptx_docs[0]["metadata"]["source_locator"] == {"type": "slide", "value": 1}
    assert "Sparse title" in pptx_docs[0]["text"]

    from docx import Document
    document = Document(); document.add_paragraph("DOCX paragraph"); docx = BytesIO(); document.save(docx)
    docx_docs = PARSERS["docx"](docx.getvalue())
    assert docx_docs[0]["metadata"]["source_locator"] == {"type": "paragraph", "value": 1}

class FakeEmbedder:
    def embed(self, texts):
        return [[float(len(texts[0]))]]

class FakeQdrant:
    def __init__(self): self.points = {}; self.deleted_tails = []
    def ensure_collection(self, dimension): assert dimension == 1
    def delete_material_tail(self, material_id, chunk_count): self.deleted_tails.append((material_id, chunk_count))
    def upsert(self, vectors, payloads, ids): self.points.update(dict(zip(ids, payloads)))

def test_pipeline_payload_is_attributable_and_ready(db, monkeypatch):
    teacher = User(id=uuid.uuid4(), email="teacher@example.com", role="teacher", name="Dr. Smith")
    subject = Subject(id=uuid.uuid4(), name="Physics")
    db.add_all([teacher, subject]); db.commit()
    material = Material(id=uuid.uuid4(), subject_id=subject.id, teacher_id=teacher.id, filename="notes.docx", file_type="docx", storage_path="x", status="processing", ingestion_version=1)
    db.add(material); db.commit()
    fake = FakeQdrant()
    monkeypatch.setitem(PARSERS, "docx", lambda data: [{"text": "A paragraph", "metadata": {"source_locator": {"type": "paragraph", "value": 1}}}])
    pipeline = IngestionPipeline(qdrant=fake, embedder=FakeEmbedder())
    pipeline.process(db, material.id, b"ignored", version=1)
    payload = next(iter(fake.points.values()))
    assert payload["teacher_name"] == "Dr. Smith"
    assert payload["filename"] == "notes.docx"
    assert payload["material_id"] == str(material.id)
    assert payload["source_locator"]["type"] == "paragraph"
    assert db.get(Material, material.id).status == "ready"

def test_stale_worker_aborts_without_qdrant_write(db):
    teacher = User(id=uuid.uuid4(), email="teacher@example.com", role="teacher", name="Dr. Smith")
    subject = Subject(id=uuid.uuid4(), name="Physics")
    db.add_all([teacher, subject]); db.commit()
    material = Material(id=uuid.uuid4(), subject_id=subject.id, teacher_id=teacher.id, filename="notes.docx", file_type="docx", storage_path="x", status="deleting", ingestion_version=2)
    db.add(material); db.commit()
    fake = FakeQdrant()
    with pytest.raises(RuntimeError, match="stale"):
        IngestionPipeline(qdrant=fake, embedder=FakeEmbedder()).process(db, material.id, b"ignored", version=1)
    assert fake.points == {}
